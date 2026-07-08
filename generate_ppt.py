from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Add 4 blank slides
for i in range(4):
    slide_layout = prs.slide_layouts[6] # 6 is completely blank
    prs.slides.add_slide(slide_layout)

# Helper function to add a title and content slide
def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if i == 0:
            p.space_before = Pt(10)

# Slide 5: Methodology
methodology_points = [
    "Event-Driven Architecture: Captures micro-events (clicks, cart adds) in real-time.",
    "Data Ingestion: Next.js frontend sends JSON payloads to FastAPI.",
    "Stream Processing: Apache Kafka acts as a shock-absorber for massive traffic.",
    "Identity Resolution: Background consumer matches anonymous cookies to real user profiles.",
    "Dynamic Profiling: Mathematical user interest scores update instantly in MySQL."
]
add_slide("Methodology: Real-Time Data Flow", methodology_points)

# Slide 6: Tools & Technologies
tech_points = [
    "Frontend: Next.js & React (Interactive E-Commerce Storefront)",
    "Commercial Backend: Django (Inventory & User Authentication)",
    "AI Engine: FastAPI (Lightning-fast Python AI serving)",
    "Event Streaming: Apache Kafka & Zookeeper (High-volume traffic queue)",
    "Database: MySQL 8.0 (Permanent Identity & Profile storage)",
    "Cache: Redis (Instant recommendation retrieval)",
    "Infrastructure: Docker (Containerized Microservices architecture)"
]
add_slide("Tools & Technologies", tech_points)

# Slide 7: Machine Learning Algorithms
ml_points = [
    "Hybrid Recommendation Engine combining three distinct approaches:",
    "1. Collaborative Filtering (ALS): Analyzes patterns (Users who bought X also bought Y).",
    "2. Content-Based Matching (BM25): Matches actual text descriptions and product metadata.",
    "3. Semantic Embeddings: Understands the real meaning of words (e.g. 'sneakers' vs 'shoes').",
    "Identity Graphing: Merges anonymous session data with logged-in user databases."
]
add_slide("Machine Learning Algorithms", ml_points)

# Slide 8: Results & Demo
results_points = [
    "Real-Time Processing: Events are captured and processed by the AI in under 10 milliseconds.",
    "Enterprise Scalability: Kafka prevents MySQL database crashes during extremely high traffic.",
    "Cold Start Solution: The 5-layer profile instantly targets brand-new users on their first click.",
    "[ DEMONSTRATION ]: (Be prepared to show the SmartShop Frontend and the Admin Dashboard)"
]
add_slide("Results & Demo", results_points)

# Slide 9: Conclusion
conclusion_points = [
    "The Contextual Product Recommender Platform (CPRP) successfully demonstrates an enterprise-grade AI architecture.",
    "By moving away from slow direct-database writes and utilizing Apache Kafka, the system achieves massive horizontal scalability.",
    "The Hybrid ML Model ensures high accuracy, solving traditional e-commerce recommendation problems like catalog decay.",
    "The project is fully containerized and ready for real-world deployment in a production environment."
]
add_slide("Conclusion", conclusion_points)

prs.save("cprp_presentation.pptx")
print("Successfully generated cprp_presentation.pptx!")
